"""
==============================================================================
GÉNÉRATEUR AVANCÉ DE PRÉSENTATIONS POWERPOINT - MODÈLES MINSANTE
==============================================================================
Module pour générer automatiquement des présentations PowerPoint avancées
avec analyse approfondie et visualisations professionnelles.

NOUVEAUX MODÈLES:
- Modèle ORIGINAL : 7 slides classiques MINSANTE (via pptx_generator_minsante.py)
- Modèle A (Amélioré) : Conserve les 7 slides MINSANTE + ajoute 9 slides analytiques = 16 SLIDES
- Modèle B (Nouvelle Version) : Design moderne complet avec 12 slides

Auteur: Fred - AIMS Cameroon / MINSANTE
Date: Décembre 2025
Version: 2.2 - Correction chemin Windows + Amélioration compatibilité multi-plateforme
==============================================================================
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import io
import tempfile
from pathlib import Path
from datetime import datetime

# Import depuis la nouvelle architecture
from config import settings


class PowerPointGeneratorAdvanced:
    """
    Générateur avancé de présentations PowerPoint avec 2 modèles.
    Utilise directement les vraies données du système.
    """
    
    # Palette de couleurs MINSANTE/Cameroun
    COLORS = {
        'vert': RGBColor(0, 122, 51),           # Vert Cameroun
        'jaune': RGBColor(255, 215, 0),         # Jaune Cameroun
        'rouge': RGBColor(206, 17, 38),         # Rouge Cameroun
        'primary': RGBColor(46, 134, 171),      # Bleu principal
        'secondary': RGBColor(162, 59, 114),    # Violet
        'accent1': RGBColor(241, 143, 1),       # Orange
        'accent2': RGBColor(199, 62, 29),       # Rouge-orange
        'accent3': RGBColor(106, 76, 147),      # Violet foncé
        'success': RGBColor(78, 205, 196),      # Turquoise
        'warning': RGBColor(255, 107, 107),     # Rouge clair
        'background': RGBColor(248, 249, 250),  # Gris très clair
        'white': RGBColor(255, 255, 255),
        'dark': RGBColor(33, 37, 41),
        'gray': RGBColor(108, 117, 125),
        'lightgray': RGBColor(233, 236, 239)
    }
    
    def __init__(self, modele="B", prs_base=None):
        """
        Initialise le générateur.
        
        Args:
            modele (str): "A" pour amélioration ou "B" pour nouvelle version
            prs_base (Presentation): Présentation de base pour le Modèle A (optionnel)
        """
        if prs_base:
            # Modèle A : On part d'une présentation existante
            self.prs = prs_base
        else:
            # Modèle B : Nouvelle présentation
            self.prs = Presentation()
            self.prs.slide_width = Inches(10)
            self.prs.slide_height = Inches(7.5)
        
        self.modele = modele
        
        # ✅ CORRECTION : Utiliser le répertoire temporaire système (Windows + Linux)
        self.charts_dir = Path(tempfile.gettempdir()) / "pptx_charts"
        self.charts_dir.mkdir(parents=True, exist_ok=True)
        
        print(f"📁 Répertoire graphiques : {self.charts_dir}")
    
    # =========================================================================
    # GÉNÉRATION DES GRAPHIQUES MATPLOTLIB
    # =========================================================================
    
    def _generer_graphique_tendances(self, df_appels, output_path):
        """Génère le graphique de tendances des appels quotidiens."""
        plt.style.use('seaborn-v0_8-darkgrid')
        fig, ax = plt.subplots(figsize=(14, 6))
        
        df = df_appels.sort_values('DATE')
        ax.plot(df['DATE'], df['TOTAL_APPELS_JOUR'], linewidth=2, 
                color='#2E86AB', label='Appels totaux')
        ax.fill_between(df['DATE'], df['TOTAL_APPELS_JOUR'], alpha=0.3, color='#2E86AB')
        
        # Ligne de moyenne
        moyenne = df['TOTAL_APPELS_JOUR'].mean()
        ax.axhline(y=moyenne, color='red', linestyle='--', linewidth=2, 
                   label=f'Moyenne: {int(moyenne)} appels/jour')
        
        ax.set_title('Évolution Quotidienne des Appels au 1510', 
                     fontsize=16, fontweight='bold', pad=20)
        ax.set_xlabel('Date', fontsize=12, fontweight='bold')
        ax.set_ylabel('Nombre d\'appels', fontsize=12, fontweight='bold')
        ax.legend(loc='upper left', fontsize=10)
        ax.grid(True, alpha=0.3)
        
        plt.xticks(rotation=45)
        plt.tight_layout()
        plt.savefig(output_path, dpi=300, bbox_inches='tight')
        plt.close()
    
    def _generer_graphique_repartition_motifs(self, df_hebdo, output_path):
        """Génère le graphique de répartition des motifs par semaine."""
        # Utiliser _SEMAINE pour df_hebdo
        motifs_cols = []
        for cat in settings.CATEGORIES_APPELS:
            col_semaine = cat.replace('_JOUR', '_SEMAINE')
            if col_semaine in df_hebdo.columns:
                motifs_cols.append(col_semaine)
        
        if not motifs_cols:
            print("⚠️ Aucune colonne de motifs trouvée dans df_hebdo")
            # Créer un graphique vide
            fig, ax = plt.subplots(figsize=(16, 8))
            ax.text(0.5, 0.5, 'Données indisponibles', ha='center', va='center', fontsize=20)
            ax.set_xticks([])
            ax.set_yticks([])
            plt.tight_layout()
            plt.savefig(output_path, dpi=300, bbox_inches='tight')
            plt.close()
            return
        
        # Prendre les 8 premières catégories disponibles
        motifs_cols = motifs_cols[:8]
        
        # Agréger par semaine
        df_weekly = df_hebdo[['Semaine épidémiologique'] + motifs_cols].copy()
        df_weekly = df_weekly.set_index('Semaine épidémiologique')
        
        fig, ax = plt.subplots(figsize=(16, 8))
        df_weekly.plot(kind='bar', stacked=True, ax=ax, width=0.8)
        
        ax.set_title('Répartition des Motifs d\'Appel par Semaine Épidémiologique',
                     fontsize=16, fontweight='bold', pad=20)
        ax.set_xlabel('Semaine Épidémiologique', fontsize=12, fontweight='bold')
        ax.set_ylabel('Nombre d\'appels', fontsize=12, fontweight='bold')
        ax.legend(title='Motifs', bbox_to_anchor=(1.05, 1), loc='upper left', fontsize=9)
        ax.grid(True, alpha=0.3, axis='y')
        
        plt.xticks(rotation=90, fontsize=8)
        plt.tight_layout()
        plt.savefig(output_path, dpi=300, bbox_inches='tight')
        plt.close()
    
    def _generer_graphique_signaux_epidemiques(self, df_appels, output_path):
        """Génère le graphique de surveillance épidémiologique."""
        fig, ax = plt.subplots(figsize=(14, 6))
        
        df = df_appels.sort_values('DATE')
        
        cols_signaux = {
            'SIGNAUX_SFE_JOUR': ('Signaux SFE', '#A23B72'),
            'CAS_SUSPECTS_JOUR': ('Cas suspects', '#F18F01'),
            'GESTION_MALADIE_JOUR': ('Gestion maladie', '#C73E1D')
        }
        
        for col, (label, color) in cols_signaux.items():
            if col in df.columns:
                ax.plot(df['DATE'], df[col], linewidth=2.5, marker='o',
                        markersize=4, label=label, color=color)
        
        ax.set_title('Surveillance Épidémiologique - Signaux et Cas Suspects',
                     fontsize=16, fontweight='bold', pad=20)
        ax.set_xlabel('Date', fontsize=12, fontweight='bold')
        ax.set_ylabel('Nombre de signaux/cas', fontsize=12, fontweight='bold')
        ax.legend(loc='upper left', fontsize=11)
        ax.grid(True, alpha=0.3)
        
        plt.xticks(rotation=45)
        plt.tight_layout()
        plt.savefig(output_path, dpi=300, bbox_inches='tight')
        plt.close()
    
    def _generer_graphique_appels_indesirables(self, df_hebdo, output_path):
        """Génère le graphique des appels indésirables."""
        cols_indesirables = []
        if 'FARCES_SEMAINE' in df_hebdo.columns:
            cols_indesirables.append('FARCES_SEMAINE')
        if 'HARCELEMENTS_SEMAINE' in df_hebdo.columns:
            cols_indesirables.append('HARCELEMENTS_SEMAINE')
        
        if not cols_indesirables:
            print("⚠️ Colonnes FARCES_SEMAINE et HARCELEMENTS_SEMAINE introuvables")
            fig, ax = plt.subplots(figsize=(16, 7))
            ax.text(0.5, 0.5, 'Données indisponibles\nColonnes FARCES et HARCELEMENTS absentes',
                    ha='center', va='center', fontsize=16, color='gray')
            ax.set_xticks([])
            ax.set_yticks([])
            plt.tight_layout()
            plt.savefig(output_path, dpi=300, bbox_inches='tight')
            plt.close()
            return
        
        df_weekly = df_hebdo.groupby('Semaine épidémiologique')[cols_indesirables].sum()
        
        fig, ax = plt.subplots(figsize=(16, 7))
        x = np.arange(len(df_weekly))
        width = 0.35
        
        if 'FARCES_SEMAINE' in cols_indesirables:
            bars1 = ax.bar(x - width/2 if len(cols_indesirables) > 1 else x, 
                          df_weekly.get('FARCES_SEMAINE', [0]*len(df_weekly)), width,
                          label='Farces', color='#FF6B6B', alpha=0.8)
        
        if 'HARCELEMENTS_SEMAINE' in cols_indesirables:
            bars2 = ax.bar(x + width/2 if len(cols_indesirables) > 1 else x, 
                          df_weekly.get('HARCELEMENTS_SEMAINE', [0]*len(df_weekly)), width,
                          label='Harcèlements', color='#4ECDC4', alpha=0.8)
        
        ax.set_title('Appels Indésirables par Semaine Épidémiologique',
                     fontsize=16, fontweight='bold', pad=20)
        ax.set_xlabel('Semaine Épidémiologique', fontsize=12, fontweight='bold')
        ax.set_ylabel('Nombre d\'appels', fontsize=12, fontweight='bold')
        ax.set_xticks(x)
        ax.set_xticklabels(df_weekly.index, rotation=90, fontsize=8)
        ax.legend(fontsize=11)
        ax.grid(True, alpha=0.3, axis='y')
        
        plt.tight_layout()
        plt.savefig(output_path, dpi=300, bbox_inches='tight')
        plt.close()
    
    def _generer_heatmap(self, df_appels, output_path):
        """Génère la heatmap d'intensité."""
        df = df_appels.copy()
        df['jour_semaine'] = pd.to_datetime(df['DATE']).dt.day_name()
        
        pivot_data = df.pivot_table(
            values='TOTAL_APPELS_JOUR',
            index='Semaine épidémiologique',
            columns='jour_semaine',
            aggfunc='mean'
        )
        
        jours_ordre = ['Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday', 'Saturday', 'Sunday']
        pivot_data = pivot_data.reindex(columns=jours_ordre, fill_value=0)
        
        fig, ax = plt.subplots(figsize=(12, 16))
        sns.heatmap(pivot_data, annot=False, cmap='YlOrRd',
                    cbar_kws={'label': 'Nombre d\'appels'},
                    linewidths=0.5, ax=ax)
        
        ax.set_title('Heatmap - Intensité des Appels par Semaine et Jour',
                     fontsize=16, fontweight='bold', pad=20)
        ax.set_xlabel('Jour de la semaine', fontsize=12, fontweight='bold')
        ax.set_ylabel('Semaine Épidémiologique', fontsize=12, fontweight='bold')
        
        jours_fr = ['Lundi', 'Mardi', 'Mercredi', 'Jeudi', 'Vendredi', 'Samedi', 'Dimanche']
        ax.set_xticklabels(jours_fr, rotation=45, ha='right')
        plt.yticks(fontsize=7)
        
        plt.tight_layout()
        plt.savefig(output_path, dpi=300, bbox_inches='tight')
        plt.close()
    
    def _generer_distribution_thematique(self, df_appels, output_path):
        """Génère le graphique de distribution thématique."""
        themes = {}
        
        # Utiliser les regroupements du config
        for groupe, categories in settings.REGROUPEMENTS.items():
            categories_existantes = [cat for cat in categories if cat in df_appels.columns]
            if categories_existantes:
                total = df_appels[categories_existantes].sum().sum()
                if total > 0:
                    themes[groupe] = total
        
        if not themes:
            print("⚠️ Aucun thème trouvé")
            fig, ax = plt.subplots(figsize=(12, 8))
            ax.text(0.5, 0.5, 'Données indisponibles', ha='center', va='center', fontsize=20)
            ax.set_xticks([])
            ax.set_yticks([])
            plt.tight_layout()
            plt.savefig(output_path, dpi=300, bbox_inches='tight')
            plt.close()
            return
        
        fig, ax = plt.subplots(figsize=(12, 8))
        colors = sns.color_palette("Set2", len(themes))
        wedges, texts, autotexts = ax.pie(
            themes.values(), labels=themes.keys(), autopct='%1.1f%%',
            colors=colors, startangle=90, textprops={'fontsize': 11}
        )
        
        for autotext in autotexts:
            autotext.set_color('white')
            autotext.set_fontweight('bold')
            autotext.set_fontsize(10)
        
        ax.set_title('Distribution Thématique des Appels',
                     fontsize=16, fontweight='bold', pad=20)
        
        plt.tight_layout()
        plt.savefig(output_path, dpi=300, bbox_inches='tight')
        plt.close()
    
    def _generer_kpi_cards(self, df_appels, output_path):
        """Génère les cartes KPI avec sparklines."""
        fig, axes = plt.subplots(2, 3, figsize=(16, 10))
        axes = axes.flatten()
        
        kpis = [
            ('Total Appels', 'TOTAL_APPELS_JOUR', '#2E86AB'),
            ('Urgences', 'URGENCE_MEDICALE_JOUR', '#A23B72'),
            ('Cas Suspects', 'CAS_SUSPECTS_JOUR', '#F18F01'),
            ('Signaux SFE', 'SIGNAUX_SFE_JOUR', '#C73E1D'),
            ('Rumeurs', 'RUMEURS_JOUR', '#6A4C93'),
            ('Farces', 'FARCES_JOUR', '#FF6B6B')
        ]
        
        df = df_appels.sort_values('DATE')
        
        for idx, (titre, colonne, couleur) in enumerate(kpis):
            ax = axes[idx]
            
            if colonne not in df.columns:
                ax.text(0.5, 0.5, f'{titre}\n\nDonnées\nindisponibles', 
                       transform=ax.transAxes,
                       fontsize=14, ha='center', va='center', color='gray')
                ax.set_xticks([])
                ax.set_yticks([])
                ax.spines['top'].set_visible(False)
                ax.spines['right'].set_visible(False)
                ax.spines['left'].set_visible(False)
                ax.spines['bottom'].set_visible(False)
                continue
            
            # Sparkline (derniers 60 jours)
            data_sparkline = df[colonne].tail(60)
            dates_sparkline = df['DATE'].tail(60)
            
            ax.plot(dates_sparkline, data_sparkline, color=couleur, linewidth=2)
            ax.fill_between(dates_sparkline, data_sparkline, alpha=0.3, color=couleur)
            
            # KPI values
            valeur_totale = df[colonne].sum()
            valeur_moyenne = df[colonne].mean()
            
            ax.text(0.5, 0.95, titre, transform=ax.transAxes,
                    fontsize=14, fontweight='bold', ha='center', va='top')
            ax.text(0.5, 0.75, f'{int(valeur_totale):,}', transform=ax.transAxes,
                    fontsize=20, fontweight='bold', ha='center', va='top',
                    color=couleur)
            ax.text(0.5, 0.60, f'Moy: {int(valeur_moyenne)}/jour',
                    transform=ax.transAxes, fontsize=10, ha='center',
                    va='top', color='gray')
            
            ax.set_xticks([])
            ax.set_yticks([])
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['left'].set_visible(False)
            ax.spines['bottom'].set_visible(False)
        
        fig.suptitle('Indicateurs Clés de Performance (KPIs)',
                     fontsize=18, fontweight='bold', y=0.98)
        plt.tight_layout()
        plt.savefig(output_path, dpi=300, bbox_inches='tight')
        plt.close()
    
    def _generer_comparaison_hebdomadaire(self, df_hebdo, output_path):
        """Génère le graphique de comparaison hebdomadaire."""
        cols_indicateurs = {
            'TOTAL_APPELS_SEMAINE': ('Total Appels', 1, '#2E86AB'),
            'URGENCE_MEDICALE_SEMAINE': ('Urgences (x10)', 10, '#A23B72'),
            'CAS_SUSPECTS_SEMAINE': ('Cas Suspects (x50)', 50, '#F18F01'),
            'SIGNAUX_SFE_SEMAINE': ('Signaux SFE (x20)', 20, '#C73E1D')
        }
        
        df_weekly = df_hebdo[['Semaine épidémiologique'] + 
                             [col for col in cols_indicateurs.keys() if col in df_hebdo.columns]].copy()
        
        fig, ax = plt.subplots(figsize=(16, 7))
        x = np.arange(len(df_weekly))
        width = 0.2
        
        positions = [-1.5, -0.5, 0.5, 1.5]
        for idx, (col, (label, mult, color)) in enumerate(cols_indicateurs.items()):
            if col in df_weekly.columns:
                ax.bar(x + positions[idx]*width, df_weekly[col]*mult, width,
                       label=label, color=color, alpha=0.8)
        
        ax.set_title('Évolution Hebdomadaire des Indicateurs Clés',
                     fontsize=16, fontweight='bold', pad=20)
        ax.set_xlabel('Semaine Épidémiologique', fontsize=12, fontweight='bold')
        ax.set_ylabel('Nombre d\'événements', fontsize=12, fontweight='bold')
        ax.set_xticks(x)
        ax.set_xticklabels(df_weekly['Semaine épidémiologique'],
                           rotation=90, fontsize=7)
        ax.legend(fontsize=10)
        ax.grid(True, alpha=0.3, axis='y')
        
        plt.tight_layout()
        plt.savefig(output_path, dpi=300, bbox_inches='tight')
        plt.close()
    
    # =========================================================================
    # CRÉATION DES SLIDES
    # =========================================================================
    
    def add_slide_with_header(self, title):
        """Crée une slide avec en-tête formaté."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Fond
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['background']
        
        # En-tête
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.9)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.COLORS['vert']
        header.line.color.rgb = self.COLORS['vert']
        
        # Titre
        title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15),
                                              Inches(9.4), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(24)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['white']
        title_para.alignment = PP_ALIGN.CENTER
        
        # Footer
        footer_box = slide.shapes.add_textbox(Inches(0.3), Inches(7.1),
                                               Inches(9.4), Inches(0.3))
        footer_frame = footer_box.text_frame
        footer_frame.text = "Centre d'Appels d'Urgence Sanitaire 1510 | MINSANTE Cameroun"
        footer_para = footer_frame.paragraphs[0]
        footer_para.font.size = Pt(9)
        footer_para.font.color.rgb = self.COLORS['gray']
        footer_para.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_image_slide(self, title, image_path):
        """Ajoute une slide avec image."""
        slide = self.add_slide_with_header(title)
        slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(1.1), width=Inches(9))
        return slide
    
    def add_title_slide(self, title, subtitle, period):
        """Ajoute une slide de titre."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Fond
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['vert']
        
        # Titre principal
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['white']
        title_para.alignment = PP_ALIGN.CENTER
        
        # Sous-titre
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = self.COLORS['white']
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Période
        period_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.8))
        period_frame = period_box.text_frame
        period_frame.text = f"Période: {period}"
        period_para = period_frame.paragraphs[0]
        period_para.font.size = Pt(18)
        period_para.font.color.rgb = self.COLORS['white']
        period_para.alignment = PP_ALIGN.CENTER
        
        # Footer
        footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(8), Inches(0.6))
        footer_frame = footer_box.text_frame
        footer_frame.text = "MINISTÈRE DE LA SANTÉ PUBLIQUE | RÉPUBLIQUE DU CAMEROUN"
        footer_para = footer_frame.paragraphs[0]
        footer_para.font.size = Pt(12)
        footer_para.font.color.rgb = self.COLORS['lightgray']
        footer_para.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_overview_slide(self, stats):
        """Ajoute une slide d'aperçu avec KPIs."""
        slide = self.add_slide_with_header("📋 APERÇU GÉNÉRAL - STATISTIQUES GLOBALES")
        
        # 3 cartes KPI principales
        kpi_data = [
            (f"{stats['total_appels']:,}", "Total Appels", self.COLORS['primary']),
            (f"{stats['moyenne_jour']}", "Moyenne/Jour", self.COLORS['secondary']),
            (f"{stats['jours_analyses']}", "Jours Analysés", self.COLORS['accent1'])
        ]
        
        x_positions = [0.5, 3.5, 6.5]
        for i, (value, label, color) in enumerate(kpi_data):
            self._add_kpi_card(slide, x_positions[i], 1.5, 2.8, 2, value, label, color)
        
        # Texte descriptif
        desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2.5))
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        
        desc_text = f"""📊 Classification des Appels (Période: {stats['period']})

✅ Appels Légitimes:
    • Urgences médicales: {stats['urgences']}
    • Cas suspects: {stats['cas_suspects']}
    • Signaux SFE: {stats['signaux_sfe']}
    • Rumeurs: {stats['rumeurs']}

⚠️ Appels Indésirables:
    • Farces: {stats['farces']:,}
    • Harcèlements: {stats['harcelements']:,}"""
        
        desc_frame.text = desc_text
        for para in desc_frame.paragraphs:
            para.font.size = Pt(13)
            para.space_after = Pt(6)
        
        return slide
    
    def _add_kpi_card(self, slide, x, y, width, height, value, label, color):
        """Ajoute une carte KPI à une slide."""
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()
        
        # Valeur
        val_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.3),
                                            Inches(width), Inches(0.8))
        val_frame = val_box.text_frame
        val_frame.text = value
        val_para = val_frame.paragraphs[0]
        val_para.font.size = Pt(36)
        val_para.font.bold = True
        val_para.font.color.rgb = self.COLORS['white']
        val_para.alignment = PP_ALIGN.CENTER
        
        # Label
        label_box = slide.shapes.add_textbox(Inches(x), Inches(y + height - 0.6),
                                              Inches(width), Inches(0.6))
        label_frame = label_box.text_frame
        label_frame.text = label
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(16)
        label_para.font.color.rgb = self.COLORS['white']
        label_para.alignment = PP_ALIGN.CENTER
    
    def add_conclusion_slide(self, stats):
        """Ajoute une slide de conclusions."""
        slide = self.add_slide_with_header("🎯 CONCLUSIONS & RECOMMANDATIONS")
        
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5),
                                                Inches(8.4), Inches(5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        conclusions_text = f"""📈 SYNTHÈSE DE L'ANALYSE ({stats['period']})

✅ POINTS CLÉS:
    • {stats['total_appels']:,} appels traités sur {stats['jours_analyses']} jours d'opération
    • Moyenne quotidienne: {stats['moyenne_jour']} appels/jour
    • {stats['urgences']} urgences médicales prises en charge
    • {stats['signaux_sfe']} signaux de surveillance épidémiologique détectés

⚠️ DÉFIS IDENTIFIÉS:
    • {stats['farces']:,} appels de farce ({stats['farces']*100//stats['total_appels'] if stats['total_appels'] > 0 else 0}% du total)
    • {stats['harcelements']:,} cas de harcèlement
    • Impact significatif sur les ressources opérationnelles

💡 RECOMMANDATIONS:
    1. Renforcer la sensibilisation publique sur l'usage approprié du 1510
    2. Développer un système de filtrage intelligent des appels indésirables
    3. Maintenir la vigilance épidémiologique continue
    4. Optimiser l'allocation des ressources pendant les pics identifiés
    5. Améliorer la formation des opérateurs sur la gestion des cas suspects"""
        
        content_frame.text = conclusions_text
        for para in content_frame.paragraphs:
            para.font.size = Pt(13)
            para.space_after = Pt(8)
        
        return slide
    
    def add_thank_you_slide(self):
        """Ajoute une slide de remerciements."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['vert']
        
        # MERCI
        thank_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
        thank_frame = thank_box.text_frame
        thank_frame.text = "MERCI"
        thank_para = thank_frame.paragraphs[0]
        thank_para.font.size = Pt(60)
        thank_para.font.bold = True
        thank_para.font.color.rgb = self.COLORS['white']
        thank_para.alignment = PP_ALIGN.CENTER
        
        # Sous-titre
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
        sub_frame = sub_box.text_frame
        sub_frame.text = "Pour votre attention"
        sub_para = sub_frame.paragraphs[0]
        sub_para.font.size = Pt(24)
        sub_para.font.color.rgb = self.COLORS['jaune']
        sub_para.alignment = PP_ALIGN.CENTER
        
        # Footer
        footer_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.8))
        footer_frame = footer_box.text_frame
        footer_frame.text = "Centre d'Appels d'Urgence Sanitaire - 1510\nMINSANTÉ | République du Cameroun"
        for para in footer_frame.paragraphs:
            para.font.size = Pt(14)
            para.font.color.rgb = self.COLORS['lightgray']
            para.alignment = PP_ALIGN.CENTER
        
        return slide
    
    # =========================================================================
    # FONCTION PRINCIPALE DE GÉNÉRATION
    # =========================================================================
    
    def generer_rapport_avance(self, df_appels, df_hebdo, date_rapport_str):
        """
        Génère un rapport PowerPoint avancé avec vraies données.
        
        Args:
            df_appels (pd.DataFrame): DataFrame des appels quotidiens
            df_hebdo (pd.DataFrame): DataFrame hebdomadaire
            date_rapport_str (str): Date formatée du rapport
        
        Returns:
            bytes: Fichier PowerPoint en bytes
        """
        # Calculer les statistiques
        stats = self._calculer_statistiques(df_appels, df_hebdo)
        stats['period'] = f"{df_appels['DATE'].min().strftime('%d/%m/%Y')} au {df_appels['DATE'].max().strftime('%d/%m/%Y')}"
        
        # Générer tous les graphiques
        print("📊 Génération des graphiques...")
        charts_paths = self._generer_tous_graphiques(df_appels, df_hebdo)
        
        if self.modele == "B":
            # MODÈLE B : Nouvelle version complète
            self._generer_modele_b(stats, charts_paths, date_rapport_str)
        else:
            # MODÈLE A : Amélioration (7 slides de base DÉJÀ PRÉSENTES + 9 analytiques)
            self._generer_modele_a(stats, charts_paths)
        
        # Sauvegarder
        return self._sauvegarder()
    
    def _calculer_statistiques(self, df_appels, df_hebdo):
        """Calcule toutes les statistiques nécessaires."""
        stats = {
            'total_appels': int(df_appels['TOTAL_APPELS_JOUR'].sum()),
            'moyenne_jour': int(df_appels['TOTAL_APPELS_JOUR'].mean()),
            'jours_analyses': len(df_appels),
            'urgences': int(df_appels['URGENCE_MEDICALE_JOUR'].sum()) if 'URGENCE_MEDICALE_JOUR' in df_appels.columns else 0,
            'cas_suspects': int(df_appels['CAS_SUSPECTS_JOUR'].sum()) if 'CAS_SUSPECTS_JOUR' in df_appels.columns else 0,
            'signaux_sfe': int(df_appels['SIGNAUX_SFE_JOUR'].sum()) if 'SIGNAUX_SFE_JOUR' in df_appels.columns else 0,
            'rumeurs': int(df_appels['RUMEURS_JOUR'].sum()) if 'RUMEURS_JOUR' in df_appels.columns else 0,
            'farces': int(df_appels['FARCES_JOUR'].sum()) if 'FARCES_JOUR' in df_appels.columns else 0,
            'harcelements': int(df_appels['HARCELEMENTS_JOUR'].sum()) if 'HARCELEMENTS_JOUR' in df_appels.columns else 0
        }
        return stats
    
    def _generer_tous_graphiques(self, df_appels, df_hebdo):
        """Génère tous les graphiques et retourne leurs chemins."""
        charts = {}
        
        print("  → Tendances...")
        charts['tendances'] = self.charts_dir / "01_tendance_appels.png"
        self._generer_graphique_tendances(df_appels, charts['tendances'])
        
        print("  → Répartition motifs...")
        charts['repartition'] = self.charts_dir / "02_repartition_motifs.png"
        self._generer_graphique_repartition_motifs(df_hebdo, charts['repartition'])
        
        print("  → Signaux épidémiques...")
        charts['signaux'] = self.charts_dir / "03_signaux_epidemiques.png"
        self._generer_graphique_signaux_epidemiques(df_appels, charts['signaux'])
        
        print("  → Appels indésirables...")
        charts['indesirables'] = self.charts_dir / "04_appels_indesirables.png"
        self._generer_graphique_appels_indesirables(df_hebdo, charts['indesirables'])
        
        print("  → Heatmap...")
        charts['heatmap'] = self.charts_dir / "05_heatmap_intensite.png"
        self._generer_heatmap(df_appels, charts['heatmap'])
        
        print("  → Distribution thématique...")
        charts['distribution'] = self.charts_dir / "06_distribution_thematique.png"
        self._generer_distribution_thematique(df_appels, charts['distribution'])
        
        print("  → KPI Cards...")
        charts['kpi'] = self.charts_dir / "07_kpi_cards.png"
        self._generer_kpi_cards(df_appels, charts['kpi'])
        
        print("  → Comparaison hebdomadaire...")
        charts['comparaison'] = self.charts_dir / "08_comparaison_hebdo.png"
        self._generer_comparaison_hebdomadaire(df_hebdo, charts['comparaison'])
        
        print("✅ Tous les graphiques générés")
        return charts
    
    def _generer_modele_a(self, stats, charts_paths):
        """
        Génère les slides supplémentaires pour le Modèle A.
        
        IMPORTANT: Les 7 slides de base MINSANTE sont DÉJÀ PRÉSENTES dans self.prs
        On ajoute seulement les 9 slides analytiques supplémentaires.
        
        Total final: 16 slides (7 de base + 9 analytiques)
        """
        print("📊 Modèle A : Ajout des 9 slides analytiques aux 7 slides MINSANTE existantes...")
        
        # Slides 8-16: Graphiques analytiques avancés
        self.add_image_slide("📊 TABLEAU DE BORD - INDICATEURS CLÉS", str(charts_paths['kpi']))
        self.add_image_slide("📈 TENDANCES - ÉVOLUTION SUR LA PÉRIODE", str(charts_paths['tendances']))
        self.add_image_slide("🎯 RÉPARTITION DES MOTIFS D'APPEL", str(charts_paths['repartition']))
        self.add_image_slide("🔬 SURVEILLANCE ÉPIDÉMIOLOGIQUE APPROFONDIE", str(charts_paths['signaux']))
        self.add_image_slide("🗓️ HEATMAP - INTENSITÉ PAR SEMAINE ET JOUR", str(charts_paths['heatmap']))
        self.add_image_slide("📊 ANALYSE THÉMATIQUE GLOBALE", str(charts_paths['distribution']))
        self.add_image_slide("⚠️ APPELS INDÉSIRABLES - FARCES ET HARCÈLEMENTS", str(charts_paths['indesirables']))
        self.add_image_slide("📊 ÉVOLUTION HEBDOMADAIRE MULTI-INDICATEURS", str(charts_paths['comparaison']))
        self.add_conclusion_slide(stats)
        
        print("✅ Modèle A : 16 slides générées (7 MINSANTE + 9 analytiques)")
    
    def _generer_modele_b(self, stats, charts_paths, date_rapport_str):
        """Génère toutes les slides pour le Modèle B (12 slides)."""
        print("📊 Modèle B : Génération de 12 slides modernes...")
        
        # Slide 1: Titre
        self.add_title_slide(
            "SURVEILLANCE DES APPELS",
            "Plateforme de Veille Sanitaire 1510",
            stats['period']
        )
        
        # Slide 2: Aperçu général
        self.add_overview_slide(stats)
        
        # Slides 3-10: Graphiques
        self.add_image_slide("📊 INDICATEURS CLÉS", str(charts_paths['kpi']))
        self.add_image_slide("📈 TENDANCES", str(charts_paths['tendances']))
        self.add_image_slide("🎯 RÉPARTITION MOTIFS", str(charts_paths['repartition']))
        self.add_image_slide("🔬 SURVEILLANCE", str(charts_paths['signaux']))
        self.add_image_slide("⚠️ INDÉSIRABLES", str(charts_paths['indesirables']))
        self.add_image_slide("🗓️ HEATMAP", str(charts_paths['heatmap']))
        self.add_image_slide("📊 DISTRIBUTION", str(charts_paths['distribution']))
        self.add_image_slide("📊 COMPARAISON", str(charts_paths['comparaison']))
        
        # Slide 11: Conclusions
        self.add_conclusion_slide(stats)
        
        # Slide 12: Remerciements
        self.add_thank_you_slide()
        
        print("✅ Modèle B : 12 slides générées")
    
    def _sauvegarder(self):
        """Sauvegarde la présentation."""
        pptx_stream = io.BytesIO()
        self.prs.save(pptx_stream)
        pptx_stream.seek(0)
        return pptx_stream.getvalue()


# ==============================================================================
# FONCTION D'INTERFACE SIMPLIFIÉE
# ==============================================================================

def generer_rapport_avance(df_appels, df_calendrier, semaine, modele, output_path):
    """
    Interface simplifiée pour générer un rapport avancé.
    
    Args:
        df_appels (pd.DataFrame): DataFrame des appels quotidiens
        df_calendrier (pd.DataFrame): DataFrame du calendrier
        semaine (str): Semaine épidémiologique (pour Modèle A uniquement)
        modele (str): "A" ou "B"
        output_path (str): Chemin de sortie
    
    Returns:
        str: Chemin du fichier généré
    """
    from utils.data_processor import calculer_totaux_hebdomadaires
    from utils.pptx_generator_minsante import MinsantePPTXGenerator
    
    print(f"🎯 Génération rapport avancé Modèle {modele}...")
    
    # Créer df_hebdo à partir de df_appels
    df_hebdo = calculer_totaux_hebdomadaires(df_appels)
    
    if modele == "A":
        # MODÈLE A : D'abord générer les 7 slides MINSANTE de base
        print("📋 Étape 1/2 : Génération des 7 slides MINSANTE de base...")
        
        # Créer le générateur MINSANTE
        gen_minsante = MinsantePPTXGenerator()
        
        # Calculer les données de la semaine
        from utils.data_processor import (
            calculer_totaux_semaine, 
            calculer_regroupements,
            comparer_periodes
        )
        
        totaux = calculer_totaux_semaine(df_appels, semaine)
        df_semaine = df_appels[df_appels['Semaine épidémiologique'] == semaine]
        regroupements = calculer_regroupements(df_semaine)
        
        # Préparer les données pour les graphiques
        renseignements_data = {}
        if 'RENSEIGNEMENTS' in settings.REGROUPEMENTS:
            for cat in settings.REGROUPEMENTS['RENSEIGNEMENTS']:
                if cat in df_semaine.columns:
                    val = int(df_semaine[cat].sum())
                    if val > 0:
                        label = settings.LABELS_CATEGORIES.get(cat, cat)
                        renseignements_data[label] = val
        
        assistance_data = {}
        if 'ASSISTANCES' in settings.REGROUPEMENTS:
            for cat in settings.REGROUPEMENTS['ASSISTANCES']:
                if cat in df_semaine.columns:
                    val = int(df_semaine[cat].sum())
                    if val > 0:
                        label = settings.LABELS_CATEGORIES.get(cat, cat)
                        assistance_data[label] = val
        
        signaux_data = {}
        if 'SIGNAUX' in settings.REGROUPEMENTS:
            for cat in settings.REGROUPEMENTS['SIGNAUX']:
                if cat in df_semaine.columns:
                    val = int(df_semaine[cat].sum())
                    if val > 0:
                        label = settings.LABELS_CATEGORIES.get(cat, cat)
                        signaux_data[label] = val
        
        autres_data = {'appels_sortants': 0, 'total': totaux['total']}
        
        # Générer les 7 slides de base
        date_rapport = totaux['date_fin'].strftime("%d %B %Y")
        periode = f"{totaux['date_debut'].strftime('%d')} au {totaux['date_fin'].strftime('%d %B %Y')}"
        
        gen_minsante.slide_1_titre(date_rapport)
        gen_minsante.slide_2_faits_saillants(periode, totaux['total'], 
                                              renseignements_data, assistance_data, 
                                              signaux_data, autres_data)
        
        # Slide 3: Comparaison
        try:
            semaines_disponibles = sorted(df_appels['Semaine épidémiologique'].unique())
            idx_actuelle = semaines_disponibles.index(semaine)
            if idx_actuelle > 0:
                semaine_precedente = semaines_disponibles[idx_actuelle - 1]
                df_comparaison = comparer_periodes(df_appels, [semaine_precedente, semaine])
                gen_minsante.slide_3_comparaison(semaine_precedente, semaine, df_comparaison)
        except Exception as e:
            print(f"⚠️ Erreur slide 3: {e}")
        
        # Slide 4: Évolution
        try:
            semaines = df_hebdo['Semaine épidémiologique'].tolist()
            valeurs = df_hebdo['TOTAL_APPELS_SEMAINE'].tolist()
            gen_minsante.slide_4_evolution(semaines, valeurs)
        except Exception as e:
            print(f"⚠️ Erreur slide 4: {e}")
        
        # Slide 5: Questions
        questions_list = [
            "Informations sur les centres de santé disponibles dans la région",
            "Symptômes de la fièvre typhoïde et traitement recommandé",
            "Disponibilité des vaccins contre la COVID-19",
            "Procédures pour signaler un cas suspect de maladie",
            "Numéros d'urgence pour les cas de traumatisme grave"
        ]
        gen_minsante.slide_5_questions_interet(periode, questions_list)
        
        # Slide 6: Activités
        activites_menees = [
            "Formation des opérateurs sur la gestion des appels d'urgence",
            "Mise à jour de la base de données des centres de santé",
            "Coordination avec les équipes de surveillance épidémiologique",
            "Analyse des tendances hebdomadaires des appels"
        ]
        activites_planifiees = [
            "Extension de la couverture géographique du service 1510",
            "Intégration d'un système de triage automatisé",
            "Formation continue sur les nouvelles pathologies émergentes",
            "Évaluation de la satisfaction des usagers"
        ]
        gen_minsante.slide_6_activites(activites_menees, activites_planifiees)
        
        # Slide 7: Merci
        gen_minsante.slide_7_merci()
        
        print("✅ 7 slides MINSANTE générées")
        
        # Étape 2 : Ajouter les 9 slides analytiques
        print("📋 Étape 2/2 : Ajout des 9 slides analytiques...")
        
        # Créer le générateur avancé avec la présentation de base
        generator = PowerPointGeneratorAdvanced(modele="A", prs_base=gen_minsante.prs)
        
    else:
        # MODÈLE B : Nouvelle présentation complète
        generator = PowerPointGeneratorAdvanced(modele="B")
    
    # Générer le rapport
    date_rapport = datetime.now().strftime("%d %B %Y")
    pptx_bytes = generator.generer_rapport_avance(df_appels, df_hebdo, date_rapport)
    
    # Sauvegarder
    with open(output_path, 'wb') as f:
        f.write(pptx_bytes)
    
    if modele == "A":
        print(f"✅ Rapport Modèle A généré : {output_path}")
        print(f"📊 16 SLIDES TOTALES : 7 MINSANTE + 9 analytiques")
    else:
        print(f"✅ Rapport Modèle B généré : {output_path}")
        print(f"📊 12 SLIDES modernes")
    
    return output_path


# ==============================================================================
# FIN DU MODULE
# ==============================================================================