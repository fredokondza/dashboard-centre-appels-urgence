"""
==============================================================================
GÉNÉRATEUR POWERPOINT - MODÈLE MINSANTE
==============================================================================
Basé EXACTEMENT sur Situation_Centre_Appel_27-11-2025.pptx
Toutes les positions et dimensions sont reproduites à l'identique

Auteur: Fred - AIMS Cameroon / MINSANTE  
Date: 14 Décembre 2025
Version: 7.0 FINAL - Positions exactes du template
==============================================================================
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from datetime import datetime
import pandas as pd
from pathlib import Path


class MinsantePPTXGenerator:
    """Générateur PowerPoint MINSANTE - Template 27-11-2025"""
    
    def __init__(self):
        """Initialise le générateur"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
        
        # Couleurs MINSANTE
        self.color_vert = RGBColor(0, 122, 51)
        self.color_jaune = RGBColor(255, 215, 0)
        self.color_rouge = RGBColor(206, 17, 38)
        self.color_white = RGBColor(255, 255, 255)
        self.color_dark = RGBColor(33, 37, 41)
        self.color_gray = RGBColor(108, 117, 125)
        self.color_blue = RGBColor(13, 110, 253)
        
        # Drapeau (optionnel)
        self.drapeau_path = None
        possible_paths = [
            Path("data/Flag_of_Cameroon.png"),
            Path("data/Flag_of_Cameroon.svg"),
        ]
        for p in possible_paths:
            if p.exists():
                self.drapeau_path = p
                break
    
    def slide_1_titre(self, date_rapport):
        """SLIDE 1 : Titre avec drapeau"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # [1] DRAPEAU - left=7.50", top=1.50", width=5.00", height=3.33"
        if self.drapeau_path:
            try:
                slide.shapes.add_picture(
                    str(self.drapeau_path),
                    Inches(7.5), Inches(1.5),
                    width=Inches(5.0), height=Inches(3.33)
                )
            except:
                pass
        
        # [2] BANDE VERTE - left=0.00", top=0.00", width=13.33", height=0.80"
        shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(13.33), Inches(0.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.color_vert
        shape.line.fill.background()
        
        # [3] TEXTE MINSANTE - left=0.50", top=0.15", width=8.00", height=0.50"
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(8.0), Inches(0.5))
        p = txBox.text_frame.paragraphs[0]
        p.text = "MINISTÈRE DE LA SANTÉ PUBLIQUE - RÉPUBLIQUE DU CAMEROUN"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.color_white
        p.alignment = PP_ALIGN.LEFT
        
        # [4] SITUATION - left=0.80", top=2.00", width=6.00", height=1.50"
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(6.0), Inches(1.5))
        p = txBox.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "SITUATION"
        run.font.size = Pt(60)
        run.font.bold = True
        run.font.color.rgb = self.color_vert
        
        # [5] SOUS-TITRE - left=0.80", top=3.50", width=6.50", height=2.00"
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(6.5), Inches(2.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "DU CENTRE D'APPELS\nD'URGENCE SANITAIRE\n"
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = self.color_dark
        
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = f"\nau {date_rapport}"
        run2.font.size = Pt(28)
        run2.font.bold = False
        run2.font.color.rgb = self.color_rouge
        
        # [6] PIED DE PAGE - left=0.50", top=6.80", width=12.33", height=0.50"
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.33), Inches(0.5))
        p = txBox.text_frame.paragraphs[0]
        p.text = "Centre d'Appels d'Urgence Sanitaire - Numéro d'urgence : 1510"
        p.font.size = Pt(14)
        p.font.color.rgb = self.color_gray
        p.alignment = PP_ALIGN.CENTER
    
    def slide_2_faits_saillants(self, periode, total_appels, renseignements_data, 
                                 assistance_data, signaux_data, autres_data):
        """SLIDE 2 : Faits saillants avec 3 graphiques"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # [1] BANDE VERTE
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.color_vert
        shape.line.fill.background()
        
        # [2] TITRE
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.0), Inches(0.6))
        p = txBox.text_frame.paragraphs[0]
        p.text = f"FAITS SAILLANTS DU {periode.upper()}"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.color_white
        p.alignment = PP_ALIGN.CENTER
        
        # [3] TOTAL APPELS
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(12.0), Inches(0.6))
        p = txBox.text_frame.paragraphs[0]
        p.text = f"📞 {total_appels:,} NOUVEAUX APPELS REÇUS".replace(",", " ")
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.color_rouge
        p.alignment = PP_ALIGN.CENTER
        
        # [4] STATISTIQUES
        total_rens = sum(renseignements_data.values()) if renseignements_data else 0
        total_assist = sum(assistance_data.values()) if assistance_data else 0
        total_sig = sum(signaux_data.values()) if signaux_data else 0
        
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(12.0), Inches(0.5))
        p = txBox.text_frame.paragraphs[0]
        p.text = f"🏥 {total_rens} Renseignements Santé  |  🚑 {total_assist} Assistances Médicales  |  📡 {total_sig} Signaux de Surveillance"
        p.font.size = Pt(16)
        p.font.color.rgb = self.color_dark
        p.alignment = PP_ALIGN.CENTER
        
        # [5] LIGNE JAUNE
        shape = slide.shapes.add_shape(1, Inches(1.5), Inches(2.6), Inches(10.33), Inches(0.02))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.color_jaune
        shape.line.fill.background()
        
        # [6] TITRE GRAPHIQUE 1
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(3.8), Inches(0.25))
        p = txBox.text_frame.paragraphs[0]
        p.text = "🏥 Renseignements Santé"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.color_vert
        p.alignment = PP_ALIGN.CENTER
        
        # [7] GRAPHIQUE 1
        self._add_pie_chart(slide, renseignements_data, 
                           Inches(0.8), Inches(3.0), Inches(3.8), Inches(3.8))
        
        # [8] TITRE GRAPHIQUE 2
        txBox = slide.shapes.add_textbox(Inches(4.8), Inches(2.7), Inches(3.8), Inches(0.25))
        p = txBox.text_frame.paragraphs[0]
        p.text = "🚑 Assistances Médicales"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.color_rouge
        p.alignment = PP_ALIGN.CENTER
        
        # [9] GRAPHIQUE 2
        self._add_pie_chart(slide, assistance_data,
                           Inches(4.8), Inches(3.0), Inches(3.8), Inches(3.8))
        
        # [10] TITRE GRAPHIQUE 3
        txBox = slide.shapes.add_textbox(Inches(8.8), Inches(2.7), Inches(3.8), Inches(0.25))
        p = txBox.text_frame.paragraphs[0]
        p.text = "📡 Signaux de Surveillance"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.color_blue
        p.alignment = PP_ALIGN.CENTER
        
        # [11] GRAPHIQUE 3
        self._add_pie_chart(slide, signaux_data,
                           Inches(8.8), Inches(3.0), Inches(3.8), Inches(3.8))
        
        # [12] PIED DE PAGE
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(12.0), Inches(0.3))
        p = txBox.text_frame.paragraphs[0]
        p.text = f"📞 {autres_data.get('appels_sortants', 0)} appel(s) sortant(s) émis  |  ⚠️ {autres_data.get('total', 0)} autres appels"
        p.font.size = Pt(14)
        p.font.color.rgb = self.color_gray
        p.alignment = PP_ALIGN.CENTER
    
    def _add_pie_chart(self, slide, data_dict, left, top, width, height):
        """Ajoute un graphique camembert"""
        if not data_dict or sum(data_dict.values()) == 0:
            data_dict = {"Aucune donnée": 1}
        
        chart_data = CategoryChartData()
        chart_data.categories = list(data_dict.keys())
        chart_data.add_series('Valeurs', list(data_dict.values()))
        
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, left, top, width, height, chart_data
        )
        
        chart = graphic_frame.chart
        chart.has_title = False
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(9)
        
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(10)
        data_labels.font.bold = True
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    
    def slide_3_comparaison(self, semaine1, semaine2, df_comparaison):
        """SLIDE 3 : Tableau de comparaison"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # [1] BANDE VERTE
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.color_vert
        shape.line.fill.background()
        
        # [2] TITRE
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.0), Inches(0.7))
        p = txBox.text_frame.paragraphs[0]
        p.text = "COMPARAISON DES APPELS"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.color_white
        p.alignment = PP_ALIGN.CENTER
        
        # [3] SOUS-TITRE
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.33), Inches(0.4))
        p = txBox.text_frame.paragraphs[0]
        p.text = f"📊 {semaine1} vs {semaine2}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.color_rouge
        p.alignment = PP_ALIGN.CENTER
        
        # [4] TABLEAU
        rows = len(df_comparaison) + 1
        cols = len(df_comparaison.columns)
        
        table = slide.shapes.add_table(
            rows, cols, Inches(1.5), Inches(1.9), Inches(10.32), Inches(5.0)
        ).table
        
        if cols == 3:
            table.columns[0].width = Inches(6.0)
            table.columns[1].width = Inches(2.16)
            table.columns[2].width = Inches(2.16)
        
        # En-têtes
        for col_idx, col_name in enumerate(df_comparaison.columns):
            cell = table.cell(0, col_idx)
            cell.text = str(col_name)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.color_vert
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.color_white
            p.alignment = PP_ALIGN.CENTER
        
        # Données
        for row_idx, row in df_comparaison.iterrows():
            for col_idx, col_name in enumerate(df_comparaison.columns):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(row[col_name])
                cell.margin_top = Inches(0.05)
                cell.margin_bottom = Inches(0.05)
                
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(248, 249, 250)
                
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(11)
                p.font.color.rgb = self.color_dark
                
                if col_idx == 0:
                    p.font.bold = True
                    p.alignment = PP_ALIGN.LEFT
                else:
                    p.alignment = PP_ALIGN.CENTER
    
    def slide_4_evolution(self, semaines, valeurs, titre_periode):
        """SLIDE 4 : Graphique d'évolution"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # [1] BANDE VERTE
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.color_vert
        shape.line.fill.background()
        
        # [2] TITRE
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.0), Inches(0.7))
        p = txBox.text_frame.paragraphs[0]
        p.text = "ÉVOLUTION DES APPELS"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.color_white
        p.alignment = PP_ALIGN.CENTER
        
        # [3] SOUS-TITRE
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.33), Inches(0.4))
        p = txBox.text_frame.paragraphs[0]
        p.text = f"📈 Période : {titre_periode}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.color_rouge
        p.alignment = PP_ALIGN.CENTER
        
        # Trier semaines
        def extraire_numero(s):
            try:
                return int(s.split('_')[0][1:])
            except:
                return 0
        
        couples = list(zip(semaines, valeurs))
        couples_tries = sorted(couples, key=lambda x: extraire_numero(x[0]))
        semaines_triees = [c[0] for c in couples_tries]
        valeurs_triees = [c[1] for c in couples_tries]
        
        # [4] GRAPHIQUE
        chart_data = CategoryChartData()
        chart_data.categories = semaines_triees
        chart_data.add_series('Appels', valeurs_triees)
        
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.8), Inches(1.9), Inches(11.73), Inches(5.2),
            chart_data
        )
        
        chart = graphic_frame.chart
        chart.has_title = False
        chart.has_legend = False
        
        series = chart.series[0]
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = self.color_vert
        
        chart.value_axis.has_major_gridlines = True
        chart.category_axis.tick_labels.font.size = Pt(9)
    
    def slide_5_questions(self, periode, questions_list):
        """SLIDE 5 : Questions d'intérêt"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # [1] BANDE VERTE
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.color_vert
        shape.line.fill.background()
        
        # [2] TITRE
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.0), Inches(0.7))
        p = txBox.text_frame.paragraphs[0]
        p.text = "QUESTIONS D'INTÉRÊT POSÉES AU 1510"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.color_white
        p.alignment = PP_ALIGN.CENTER
        
        # [3] SOUS-TITRE
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.33), Inches(0.4))
        p = txBox.text_frame.paragraphs[0]
        p.text = f"📅 Période : {periode}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.color_rouge
        p.alignment = PP_ALIGN.CENTER
        
        # [4] QUESTIONS
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.33), Inches(4.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for i, question in enumerate(questions_list, 1):
            if i > 1:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
            
            p.text = f"{i}. {question}"
            p.font.size = Pt(15)
            p.font.color.rgb = self.color_dark
            p.space_before = Pt(12)
            p.space_after = Pt(12)
            p.line_spacing = 1.3
    
    def slide_6_activites(self, activites_menees, activites_planifiees):
        """SLIDE 6 : Activités"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # [1] BANDE VERTE
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.color_vert
        shape.line.fill.background()
        
        # [2] TITRE
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.0), Inches(0.6))
        p = txBox.text_frame.paragraphs[0]
        p.text = "ACTIVITÉS MENÉES ET PLANIFIÉES"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.color_white
        p.alignment = PP_ALIGN.CENTER
        
        # [3] TABLEAU
        table = slide.shapes.add_table(
            2, 2, Inches(1.0), Inches(1.5), Inches(11.33), Inches(5.5)
        ).table
        
        table.columns[0].width = Inches(5.665)
        table.columns[1].width = Inches(5.665)
        
        # En-têtes
        headers = ["✅ ACTIVITÉS MENÉES", "📋 ACTIVITÉS PLANIFIÉES"]
        for col_idx, header_text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.color_vert
            cell.margin_top = Inches(0.1)
            cell.margin_bottom = Inches(0.1)
            
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.color_white
            p.alignment = PP_ALIGN.CENTER
        
        # Activités menées
        cell = table.cell(1, 0)
        cell.margin_left = Inches(0.15)
        cell.margin_right = Inches(0.15)
        cell.margin_top = Inches(0.15)
        tf = cell.text_frame
        tf.clear()
        
        for i, activite in enumerate(activites_menees):
            if i > 0:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
            p.text = f"• {activite}"
            p.font.size = Pt(13)
            p.font.color.rgb = self.color_dark
            p.space_before = Pt(8)
            p.space_after = Pt(8)
        
        # Activités planifiées
        cell = table.cell(1, 1)
        cell.margin_left = Inches(0.15)
        cell.margin_right = Inches(0.15)
        cell.margin_top = Inches(0.15)
        tf = cell.text_frame
        tf.clear()
        
        for i, activite in enumerate(activites_planifiees):
            if i > 0:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
            p.text = f"• {activite}"
            p.font.size = Pt(13)
            p.font.color.rgb = self.color_dark
            p.space_before = Pt(8)
            p.space_after = Pt(8)
    
    def slide_7_merci(self):
        """SLIDE 7 : Merci"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # [1] FOND VERT
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.color_vert
        shape.line.fill.background()
        
        # [2] TEXTE MERCI
        txBox = slide.shapes.add_textbox(Inches(2.0), Inches(2.0), Inches(9.33), Inches(2.0))
        tf = txBox.text_frame
        tf.text = "MERCI"
        
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(90)
        p.font.bold = True
        p.font.color.rgb = self.color_white
        
        # [3] SOUS-TEXTE
        txBox = slide.shapes.add_textbox(Inches(2.0), Inches(4.2), Inches(9.33), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = "Pour votre attention"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(28)
        p.font.color.rgb = self.color_jaune
        
        p2 = tf.add_paragraph()
        p2.text = "\nCentre d'Appels d'Urgence Sanitaire - 1510"
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(18)
        p2.font.color.rgb = self.color_white
    
    def sauvegarder(self, output_path):
        """Sauvegarde la présentation"""
        self.prs.save(output_path)


# ==============================================================================
# FONCTION WRAPPER
# ==============================================================================

def generer_rapport_minsante(df_appels, df_calendrier, semaine, output_path):
    """Génère un rapport PowerPoint MINSANTE"""
    from utils.data_processor import (
        calculer_totaux_semaine,
        calculer_totaux_hebdomadaires,
        comparer_periodes
    )
    from config import settings
    
    # Calculer totaux
    totaux = calculer_totaux_semaine(df_appels, semaine)
    df_semaine = df_appels[df_appels['Semaine épidémiologique'] == semaine]
    
    # Préparer données graphiques
    # ✅ CORRECTION : Utiliser les VRAIES clés de settings.REGROUPEMENTS
    
    # Graphique 1 : Renseignements Santé
    renseignements_data = {}
    if 'Renseignements Santé' in settings.REGROUPEMENTS:
        for cat in settings.REGROUPEMENTS['Renseignements Santé']:
            if cat in df_semaine.columns:
                val = int(df_semaine[cat].sum())
                if val > 0:
                    label = settings.LABELS_CATEGORIES.get(cat, cat)
                    renseignements_data[label] = val
    
    # Graphique 2 : Assistances Médicales
    assistance_data = {}
    if 'Assistances Médicales' in settings.REGROUPEMENTS:
        for cat in settings.REGROUPEMENTS['Assistances Médicales']:
            if cat in df_semaine.columns:
                val = int(df_semaine[cat].sum())
                if val > 0:
                    label = settings.LABELS_CATEGORIES.get(cat, cat)
                    assistance_data[label] = val
    
    # Graphique 3 : Signaux
    signaux_data = {}
    if 'Signaux' in settings.REGROUPEMENTS:
        for cat in settings.REGROUPEMENTS['Signaux']:
            if cat in df_semaine.columns:
                val = int(df_semaine[cat].sum())
                if val > 0:
                    label = settings.LABELS_CATEGORIES.get(cat, cat)
                    signaux_data[label] = val
    
    # Créer générateur
    gen = MinsantePPTXGenerator()
    
    # SLIDE 1
    date_rapport = totaux['date_fin'].strftime("%d %B %Y")
    gen.slide_1_titre(date_rapport)
    
    # SLIDE 2
    periode = f"{totaux['date_debut'].strftime('%d')} au {totaux['date_fin'].strftime('%d %B %Y')}"
    gen.slide_2_faits_saillants(
        periode, totaux['total'],
        renseignements_data, assistance_data, signaux_data,
        {'appels_sortants': 0, 'total': totaux['total']}
    )
    
    # SLIDE 3 : Comparaison (TABLEAU)
    try:
        semaines_disponibles = sorted(df_appels['Semaine épidémiologique'].unique())
        idx = semaines_disponibles.index(semaine)
        if idx > 0:
            semaine_precedente = semaines_disponibles[idx - 1]
            df_comp = comparer_periodes(df_appels, [semaine_precedente, semaine])
            gen.slide_3_comparaison(semaine_precedente, semaine, df_comp)
        else:
            # Créer un DataFrame vide si pas de semaine précédente
            import pandas as pd
            df_vide = pd.DataFrame({
                'Catégorie': ['Pas de données'],
                semaine: [0]
            })
            gen.slide_3_comparaison(semaine, semaine, df_vide)
    except Exception as e:
        print(f"⚠️ Erreur slide 3: {e}")
        import pandas as pd
        df_vide = pd.DataFrame({
            'Catégorie': ['Erreur'],
            semaine: [0]
        })
        gen.slide_3_comparaison(semaine, semaine, df_vide)
    
    # SLIDE 4 : Évolution (GRAPHIQUE) - TOUJOURS avec TOUTES les semaines
    try:
        # ✅ Utiliser df_appels COMPLET pour avoir l'évolution de S5_2025 à S48_2025
        df_hebdo = calculer_totaux_hebdomadaires(df_appels)
        
        # Fonction de tri
        def extraire_numero(s):
            try:
                return int(s.split('_')[0][1:])
            except:
                return 0
        
        # Récupérer et trier les semaines
        semaines = df_hebdo['Semaine épidémiologique'].tolist()
        valeurs = df_hebdo['TOTAL_APPELS_SEMAINE'].tolist()
        
        couples = list(zip(semaines, valeurs))
        couples_tries = sorted(couples, key=lambda x: extraire_numero(x[0]))
        semaines_triees = [c[0] for c in couples_tries]
        valeurs_triees = [c[1] for c in couples_tries]
        
        # Titre avec première et dernière semaine
        titre = f"{semaines_triees[0]} à {semaines_triees[-1]}"
        
        print(f"📈 Graphique évolution : {len(semaines_triees)} semaines de {semaines_triees[0]} à {semaines_triees[-1]}")
        
        gen.slide_4_evolution(semaines_triees, valeurs_triees, titre)
    except Exception as e:
        print(f"⚠️ Erreur slide 4: {e}")
        # En cas d'erreur, créer graphique avec la semaine actuelle seulement
        gen.slide_4_evolution([semaine], [totaux['total']], semaine)
    
    # SLIDE 5 : Questions (TEXTE)
    questions = [
        "Qu'elle est la durée de validité d'une carte CSU ?",
        "Est-ce qu'un Diabétique peut bénéficier de la CSU ?",
        "Comment obtenir une carte CSU ?",
        "Quels sont les centres de santé agréés CSU ?",
        "Puis-je utiliser ma carte CSU dans toutes les régions ?"
    ]
    gen.slide_5_questions(periode, questions)
    
    # SLIDE 6 : Activités (TABLEAU)
    activites_menees = [
        "Formation des opérateurs sur la gestion des appels d'urgence",
        "Mise à jour de la base de données des centres de santé",
        "Coordination avec les équipes de surveillance épidémiologique"
    ]
    activites_planifiees = [
        "Extension de la couverture géographique du service 1510",
        "Intégration d'un système de triage automatisé",
        "Formation continue sur les nouvelles pathologies émergentes"
    ]
    gen.slide_6_activites(activites_menees, activites_planifiees)
    
    # SLIDE 7 : Merci
    gen.slide_7_merci()
    
    # Sauvegarder
    gen.sauvegarder(output_path)
    
    return output_path